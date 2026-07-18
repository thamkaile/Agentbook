from __future__ import annotations

import tempfile
import unittest

from io import BytesIO
from pathlib import Path

from langchain_core.documents import Document
from pptx import Presentation
from pptx.util import Inches

from backend.rag.ingestion import prepare_chunks
from backend.rag.loaders import get_mime_type, load_documents_from_bytes


PPTX_MIME_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)


class ByteLoaderTest(unittest.TestCase):
    def test_txt_bytes_extract_without_persisting_temporary_source_path(
        self,
    ) -> None:
        documents = load_documents_from_bytes(
            filename="lesson.txt",
            file_data=b"First idea.\n\nSecond idea.",
        )

        self.assertEqual(len(documents), 1)
        self.assertEqual(
            documents[0].page_content,
            "First idea.\n\nSecond idea.",
        )
        self.assertEqual(documents[0].metadata["filename"], "lesson.txt")
        self.assertNotIn("source", documents[0].metadata)

    def test_pptx_extracts_text_tables_and_groups_per_nonempty_slide(
        self,
    ) -> None:
        presentation_bytes = self._build_presentation()

        documents = load_documents_from_bytes(
            filename="plants.pptx",
            file_data=presentation_bytes,
        )

        self.assertEqual(len(documents), 2)
        self.assertEqual(
            [document.metadata["slide_number"] for document in documents],
            [1, 3],
        )
        first_slide = documents[0].page_content
        all_text = "\n".join(document.page_content for document in documents)

        self.assertIn("Photosynthesis", first_slide)
        self.assertEqual(first_slide.count("Photosynthesis"), 1)
        self.assertIn("Plants convert light energy.", first_slide)
        self.assertIn("Input", first_slide)
        self.assertIn("Light", first_slide)
        self.assertIn("Grouped detail", first_slide)
        self.assertIn("Cellular respiration", all_text)
        self.assertNotIn("PRIVATE SPEAKER NOTE", all_text)

        for document in documents:
            self.assertEqual(document.metadata["filename"], "plants.pptx")
            self.assertNotIn("source", document.metadata)

    def test_pptx_and_txt_mime_types(self) -> None:
        self.assertEqual(get_mime_type(Path("notes.txt")), "text/plain")
        self.assertEqual(
            get_mime_type(Path("lecture.pptx")),
            PPTX_MIME_TYPE,
        )

    def test_empty_text_and_corrupt_pptx_are_rejected(self) -> None:
        with self.assertRaises(ValueError):
            load_documents_from_bytes(
                filename="empty.txt",
                file_data=b"  \n\t",
            )

        with self.assertRaises((ValueError, RuntimeError)):
            load_documents_from_bytes(
                filename="broken.pptx",
                file_data=b"not an Office archive",
            )

    def test_chunk_metadata_keeps_mime_and_separate_slide_lineage(self) -> None:
        chunks = prepare_chunks(
            documents=[
                Document(
                    page_content="Slide-specific evidence.",
                    metadata={
                        "filename": "lecture.pptx",
                        "mime_type": PPTX_MIME_TYPE,
                        "slide_number": 4,
                    },
                )
            ],
            document_id=17,
            filename="lecture.pptx",
        )

        self.assertEqual(len(chunks), 1)
        metadata = chunks[0].metadata
        self.assertEqual(metadata["document_id"], 17)
        self.assertEqual(metadata["filename"], "lecture.pptx")
        self.assertEqual(metadata["mime_type"], PPTX_MIME_TYPE)
        self.assertEqual(metadata["slide_number"], 4)
        self.assertEqual(metadata["chunk_index"], 0)
        self.assertNotEqual(metadata.get("page_number"), 4)

    @staticmethod
    def _build_presentation() -> bytes:
        presentation = Presentation()
        blank_layout = presentation.slide_layouts[6]

        first = presentation.slides.add_slide(blank_layout)
        title = first.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(8), Inches(0.6)
        )
        title.text = "Photosynthesis"
        duplicate_title = first.shapes.add_textbox(
            Inches(0.5), Inches(1.0), Inches(8), Inches(0.6)
        )
        duplicate_title.text = "Photosynthesis"
        body = first.shapes.add_textbox(
            Inches(0.5), Inches(1.6), Inches(8), Inches(0.8)
        )
        body.text = "Plants convert light energy."

        table_shape = first.shapes.add_table(
            2,
            2,
            Inches(0.5),
            Inches(2.5),
            Inches(5),
            Inches(1.2),
        )
        table_shape.table.cell(0, 0).text = "Input"
        table_shape.table.cell(0, 1).text = "Output"
        table_shape.table.cell(1, 0).text = "Light"
        table_shape.table.cell(1, 1).text = "Sugar"

        group = first.shapes.add_group_shape()
        grouped_text = group.shapes.add_textbox(
            Inches(0.5), Inches(4.0), Inches(5), Inches(0.6)
        )
        grouped_text.text = "Grouped detail"

        presentation.slides.add_slide(blank_layout)

        third = presentation.slides.add_slide(blank_layout)
        third_body = third.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(8), Inches(0.8)
        )
        third_body.text = "Cellular respiration"
        third.notes_slide.notes_text_frame.text = "PRIVATE SPEAKER NOTE"

        buffer = BytesIO()
        presentation.save(buffer)
        return buffer.getvalue()


if __name__ == "__main__":
    unittest.main()
