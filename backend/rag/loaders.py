from __future__ import annotations

from io import BytesIO
import tempfile
from pathlib import Path
from typing import Any, Iterator

from langchain_community.document_loaders import PyPDFLoader, TextLoader
from langchain_core.documents import Document


SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".txt"}


def validate_file_path(file_path: Path) -> None:
    if not file_path.exists():
        raise FileNotFoundError(f"File does not exist: {file_path}")

    if not file_path.is_file():
        raise ValueError(f"Path is not a file: {file_path}")

    if file_path.suffix.lower() not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValueError(
            f"Unsupported file type: {file_path.suffix}. "
            f"Supported types: {supported}"
        )


def get_mime_type(file_path: Path) -> str:
    extension = file_path.suffix.lower()

    if extension == ".pdf":
        return "application/pdf"

    if extension == ".txt":
        return "text/plain"

    if extension == ".pptx":
        return (
            "application/vnd.openxmlformats-officedocument."
            "presentationml.presentation"
        )

    return "application/octet-stream"


def _normalized_text_key(text: str) -> str:
    """Normalize extracted text only for duplicate detection."""
    return " ".join(text.split()).casefold()


def _iter_text_frame_segments(text_frame: Any) -> Iterator[str]:
    for paragraph in text_frame.paragraphs:
        text = paragraph.text.strip()

        if text:
            yield text


def _iter_table_segments(table: Any) -> Iterator[str]:
    for row in table.rows:
        cell_texts: list[str] = []

        for cell in row.cells:
            paragraphs = list(
                _iter_text_frame_segments(cell.text_frame)
            )

            if paragraphs:
                cell_texts.append(" ".join(paragraphs))

        if cell_texts:
            yield " | ".join(cell_texts)


def _iter_shape_segments(
    shape: Any,
    *,
    group_shape_type: Any,
) -> Iterator[str]:
    """Yield readable text while ignoring non-text visuals."""
    if shape.shape_type == group_shape_type:
        for nested_shape in shape.shapes:
            yield from _iter_shape_segments(
                nested_shape,
                group_shape_type=group_shape_type,
            )

        return

    if bool(getattr(shape, "has_table", False)):
        yield from _iter_table_segments(shape.table)
        return

    if bool(getattr(shape, "has_text_frame", False)):
        yield from _iter_text_frame_segments(
            shape.text_frame
        )


def _load_presentation_documents(
    filename: str,
    file_data: bytes,
) -> list[Document]:
    """Extract text-only logical sections from modern PowerPoint slides."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        presentation = Presentation(BytesIO(file_data))

    except ImportError as error:
        raise RuntimeError(
            "PowerPoint support requires the python-pptx package."
        ) from error

    except Exception as error:
        raise ValueError(
            "The PowerPoint file could not be read. It may be "
            "corrupt, protected, or not a valid .pptx file."
        ) from error

    documents: list[Document] = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1,
    ):
        segments: list[str] = []
        seen_segments: set[str] = set()

        for shape in slide.shapes:
            for raw_segment in _iter_shape_segments(
                shape,
                group_shape_type=MSO_SHAPE_TYPE.GROUP,
            ):
                segment = raw_segment.strip()
                normalized = _normalized_text_key(segment)

                if not normalized or normalized in seen_segments:
                    continue

                seen_segments.add(normalized)
                segments.append(segment)

        if not segments:
            continue

        documents.append(
            Document(
                page_content="\n\n".join(segments),
                metadata={
                    "filename": filename,
                    "mime_type": get_mime_type(Path(filename)),
                    "slide_number": slide_number,
                },
            )
        )

    return documents


def load_documents_from_bytes(
    filename: str,
    file_data: bytes,
) -> list[Document]:
    suffix = Path(filename).suffix.lower()

    if suffix not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {suffix}")

    if suffix == ".pptx":
        documents = _load_presentation_documents(
            filename=filename,
            file_data=file_data,
        )

        if not documents:
            raise ValueError(
                "No readable slide text was extracted from the "
                "PowerPoint file. Charts, diagrams, screenshots, "
                "images, and visual equations are not interpreted."
            )

        return documents

    temporary_path: Path | None = None

    try:
        with tempfile.NamedTemporaryFile(
            mode="wb",
            suffix=suffix,
            delete=False,
        ) as temporary_file:
            temporary_file.write(file_data)
            temporary_path = Path(temporary_file.name)

        if suffix == ".pdf":
            loader = PyPDFLoader(str(temporary_path))
        else:
            loader = TextLoader(
                str(temporary_path),
                encoding="utf-8",
                autodetect_encoding=True,
            )

        try:
            documents = loader.load()
        except Exception as error:
            if suffix == ".pdf":
                message = (
                    "The PDF could not be read. It may be corrupt, "
                    "password-protected, or image-based."
                )
            else:
                message = (
                    "The text file could not be decoded as readable "
                    "text."
                )

            raise ValueError(message) from error

        valid_documents = [
            document
            for document in documents
            if document.page_content.strip()
        ]

        if not valid_documents:
            raise ValueError(
                "No readable text was extracted from the file. "
                "The PDF may be scanned or image-based."
            )

        for document in valid_documents:
            document.metadata.pop("source", None)
            document.metadata["filename"] = filename
            document.metadata["mime_type"] = get_mime_type(
                Path(filename)
            )

        return valid_documents

    finally:
        if temporary_path is not None:
            temporary_path.unlink(missing_ok=True)
